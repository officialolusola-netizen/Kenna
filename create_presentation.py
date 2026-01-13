from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color scheme (professional blue and green)
    PRIMARY_COLOR = RGBColor(0, 51, 102)  # Navy blue
    SECONDARY_COLOR = RGBColor(0, 128, 96)  # Green
    ACCENT_COLOR = RGBColor(255, 153, 0)  # Orange
    TEXT_COLOR = RGBColor(51, 51, 51)  # Dark gray

    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add colored background shape
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = PRIMARY_COLOR
        bg.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(200, 200, 200)
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Add date and presenter info
        info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
        info_frame = info_box.text_frame
        info_frame.text = "January 8, 2026"
        info_para = info_frame.paragraphs[0]
        info_para.font.size = Pt(18)
        info_para.font.color.rgb = RGBColor(200, 200, 200)
        info_para.alignment = PP_ALIGN.CENTER

    def add_content_slide(title, content_items, slide_number=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # Add title bar
        title_bar = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = PRIMARY_COLOR
        title_bar.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8.5), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(10)
            p.level = 0

            # Add bullet point
            p.text = "• " + item

        # Add slide number if provided
        if slide_number:
            num_box = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.8), Inches(0.3))
            num_frame = num_box.text_frame
            num_frame.text = str(slide_number)
            num_para = num_frame.paragraphs[0]
            num_para.font.size = Pt(12)
            num_para.font.color.rgb = TEXT_COLOR
            num_para.alignment = PP_ALIGN.RIGHT

    def add_section_slide(section_title):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add colored background
        bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = SECONDARY_COLOR
        bg.line.fill.background()

        # Add section title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = section_title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

    # Slide 1: Title Slide
    add_title_slide(
        "Credit Facilitation in Nigeria",
        "Frameworks, Challenges, and Opportunities"
    )

    # Slide 2: Overview/Agenda
    add_content_slide(
        "Presentation Overview",
        [
            "Definition and Importance of Credit Facilitation",
            "Key Players in Nigeria's Credit Ecosystem",
            "Regulatory Framework and CBN's Role",
            "Credit Instruments and Guarantee Schemes",
            "Credit Infrastructure: Bureaus and Collateral Registry",
            "Challenges Limiting Credit Access",
            "Recent Regulatory Interventions",
            "Opportunities and Way Forward"
        ],
        2
    )

    # Slide 3: What is Credit Facilitation?
    add_content_slide(
        "What is Credit Facilitation?",
        [
            "Definition: Mechanisms, instruments, and institutional frameworks enabling credit flow from lenders to borrowers",
            "Critical for economic growth and development",
            "Encompasses lending infrastructure, regulatory systems, and risk mitigation tools",
            "Reduces information asymmetry between lenders and borrowers",
            "Facilitates financial inclusion and economic opportunity"
        ],
        3
    )

    # Slide 4: Why Credit Facilitation Matters
    add_content_slide(
        "Why Credit Facilitation Matters for Nigeria",
        [
            "SMEs contribute significantly to GDP but face severe credit constraints",
            "MSMEs face an unmet credit demand of $32.2 billion (IFC Study)",
            "Only 13.9% credit bureau coverage vs 66.5% in South Africa",
            "Agriculture employs two-thirds of workforce but lacks adequate financing",
            "Credit access drives entrepreneurship, job creation, and poverty reduction"
        ],
        4
    )

    # Section: Key Players
    add_section_slide("KEY PLAYERS IN NIGERIA'S CREDIT ECOSYSTEM")

    # Slide 5: Commercial Banks
    add_content_slide(
        "Commercial Banks",
        [
            "Primary source of formal credit in Nigeria",
            "Ongoing recapitalization requirements to strengthen capital buffers",
            "Focus on corporate and high-net-worth individuals",
            "High collateral requirements limit SME access",
            "CBN's 2026 banking reforms aim to enhance credit intermediation",
            "Interest rates remain relatively high due to macroeconomic risks"
        ],
        5
    )

    # Slide 6: Development Finance Institutions
    add_content_slide(
        "Development Finance Institutions (DFIs)",
        [
            "Development Bank of Nigeria (DBN): Leading DFI for MSMEs",
            "Provided on-lending to 321,867 MSMEs (66% women-owned, 12% first-time borrowers)",
            "Bank of Industry (BOI): Industrial and manufacturing sector focus",
            "Bank of Agriculture (BOA): Agricultural value chain financing",
            "Offer longer tenure and lower rates compared to commercial banks",
            "Critical for sectors underserved by commercial banks"
        ],
        6
    )

    # Slide 7: Microfinance Banks
    add_content_slide(
        "Microfinance Banks (MFBs)",
        [
            "Target low-income individuals and micro-enterprises",
            "Digital MFBs showing rapid growth: Moniepoint reached unicorn status (2024)",
            "Over 900 licensed MFBs nationwide providing grassroots access",
            "Challenges: Limited capital base, high operating costs, portfolio quality",
            "Increasing digitalization improving reach and efficiency",
            "Critical for last-mile financial inclusion"
        ],
        7
    )

    # Slide 8: Fintech Lenders
    add_content_slide(
        "Fintech Lenders",
        [
            "Number of approved digital lenders surged to 461 (August 2025)",
            "Nigeria Mobile Lending & Fintech Market valued at $2.1 billion",
            "Leading players: Moniepoint, FairMoney (10,000+ daily loans), Carbon",
            "Leverage alternative data and AI for credit scoring",
            "Enable faster loan approvals reaching underserved segments",
            "Challenges: Rising borrower indebtedness, regulatory compliance"
        ],
        8
    )

    # Section: Regulatory Framework
    add_section_slide("REGULATORY FRAMEWORK")

    # Slide 9: CBN's Regulatory Role
    add_content_slide(
        "Central Bank of Nigeria (CBN) - Regulatory Framework",
        [
            "Primary regulator of Nigeria's financial system and credit market",
            "Banking Supervision: Capital adequacy, lending standards, risk management",
            "Monetary Policy: Interest rate setting, liquidity management",
            "Development Finance: Intervention funds and guarantee schemes",
            "Consumer Protection: Guidelines for fair lending practices",
            "Financial Inclusion mandates and targets"
        ],
        9
    )

    # Slide 10: Recent CBN Interventions (2025-2026)
    add_content_slide(
        "Recent CBN Regulatory Interventions (2025-2026)",
        [
            "Open Banking Framework: Implementation began August 2025, full roll-out early 2026",
            "Agent Banking Guidelines (October 2025): Exclusivity clauses, transaction limits",
            "Banking Recapitalization: Strengthening capital buffers for enhanced lending",
            "Guidelines for Credit Guarantee Companies: Formal regulation introduced",
            "Digital Lending Oversight: FCCPC Consumer Lending Regulations (July 2025)",
            "Enhanced credit bureau reporting mandates"
        ],
        10
    )

    # Slide 11: Credit Instruments
    add_content_slide(
        "Credit Instruments and Mechanisms",
        [
            "Term Loans: Fixed period loans for working capital and asset acquisition",
            "Overdraft Facilities: Short-term revolving credit for cash flow management",
            "Trade Finance: Letters of credit, bills discounting for international trade",
            "Lease Financing: Asset acquisition without large upfront capital",
            "Invoice Discounting: Working capital against receivables",
            "Supply Chain Finance: Financing tied to buyer-supplier relationships"
        ],
        11
    )

    # Section: Guarantee Schemes
    add_section_slide("CREDIT GUARANTEE SCHEMES")

    # Slide 12: ACGSF
    add_content_slide(
        "Agricultural Credit Guarantee Scheme Fund (ACGSF)",
        [
            "Established 1977 (Decree 20), commenced operations April 1978",
            "Managed by CBN; Funded 60% Federal Govt, 40% CBN",
            "Share capital expanded from ₦3 billion to ₦50 billion (2019)",
            "Guarantees 75% of loan amount in default net of security realized",
            "Interest Drawback Programme reduces effective borrowing rates",
            "New Board (Dec 2025) charged with eliminating collateral/location barriers"
        ],
        12
    )

    # Slide 13: Other Guarantee Schemes
    add_content_slide(
        "Other Credit Guarantee Schemes",
        [
            "SME Credit Guarantee Scheme (SMECGS): Guarantees up to 80% of SME loans",
            "National Collateral Registry Guarantee: Enables movable asset collateralization",
            "Export Expansion Grant: Supports export-oriented businesses",
            "Youth Entrepreneurship Support (YES): Targets young entrepreneurs",
            "Women's Access to Credit: Gender-focused guarantee programs",
            "Challenge: Low awareness and utilization rates among target beneficiaries"
        ],
        13
    )

    # Section: Credit Infrastructure
    add_section_slide("CREDIT INFRASTRUCTURE")

    # Slide 14: Credit Bureaus
    add_content_slide(
        "Credit Bureaus in Nigeria",
        [
            "Licensed bureaus: CRC Credit Bureau, FirstCentral, XDS, others",
            "Credit Reporting Act mandates commercial banks report to bureaus",
            "Current coverage: 13.9% (significantly below South Africa's 66.5%)",
            "Challenges: Incomplete reporting, loopholes in the system, low awareness",
            "2025 Digital Lending Regulations require biannual bureau reporting",
            "Potential: Improved data sharing could reduce default rates and lending costs"
        ],
        14
    )

    # Slide 15: National Collateral Registry
    add_content_slide(
        "National Collateral Registry (NCR)",
        [
            "Established to enable use of movable assets as collateral",
            "Addresses traditional over-reliance on land/property collateral",
            "Only 6% of borrowers have had collateral registered on their behalf",
            "Only 33% of credit providers report customer awareness of registry",
            "Some lenders still unwilling to accept movable assets as collateral",
            "Opportunity: Better utilization could significantly expand credit access"
        ],
        15
    )

    # Section: Challenges
    add_section_slide("CHALLENGES LIMITING CREDIT ACCESS")

    # Slide 16: High Interest Rates
    add_content_slide(
        "Challenge 1: High Interest Rates",
        [
            "Lending rates remain elevated due to high MPR and macroeconomic risks",
            "Inflation pressures keeping real rates volatile",
            "Risk premiums for SMEs and individual borrowers particularly high",
            "Cost of funds limits borrower ability to service loans profitably",
            "CBN projects inflation toward single digits in 2026 (potential relief)",
            "Need for targeted interventions to reduce lending costs"
        ],
        16
    )

    # Slide 17: Collateral Requirements
    add_content_slide(
        "Challenge 2: Restrictive Collateral Requirements",
        [
            "Traditional preference for land/property as collateral",
            "Many SMEs and individuals lack acceptable collateral",
            "Movable asset collateralization underutilized despite NCR existence",
            "Lender reluctance to accept alternative collateral forms",
            "Valuation and enforcement challenges for non-traditional collateral",
            "Excludes large segment of creditworthy but asset-poor borrowers"
        ],
        17
    )

    # Slide 18: Legal and Institutional Gaps
    add_content_slide(
        "Challenge 3: Legal and Institutional Gaps",
        [
            "Weak creditor rights enforcement mechanisms",
            "Lengthy and costly loan recovery processes through courts",
            "Insolvency framework inadequate for efficient asset recovery",
            "Land title registration systems fragmented and unreliable",
            "Limited capacity of specialized commercial courts",
            "Creates reluctance among lenders, increases risk premiums"
        ],
        18
    )

    # Slide 19: Information Asymmetry
    add_content_slide(
        "Challenge 4: Information Asymmetry",
        [
            "Low credit bureau coverage (13.9%) limits credit history access",
            "Incomplete reporting by lenders creates information gaps",
            "Many borrowers lack formal financial records/documentation",
            "Digital lenders facilitating 'loan stacking' and multiple defaults",
            "Alternative data sources underutilized in credit assessment",
            "Increases perceived lending risk and rejection rates"
        ],
        19
    )

    # Slide 20: Macroeconomic Instability
    add_content_slide(
        "Challenge 5: Macroeconomic Instability",
        [
            "Exchange rate volatility affects loan pricing and servicing",
            "Inflation erodes borrower purchasing power and repayment capacity",
            "Fiscal pressures limit government support for intervention programs",
            "Business environment uncertainty reduces credit demand",
            "Policy inconsistency creates planning difficulties for lenders",
            "CBN forecasts improvement: GDP growth acceleration expected in 2026"
        ],
        20
    )

    # Section: Opportunities
    add_section_slide("OPPORTUNITIES AND WAY FORWARD")

    # Slide 21: Digital Innovation
    add_content_slide(
        "Opportunity 1: Digital Financial Innovation",
        [
            "Open Banking Framework enables better credit assessment via data sharing",
            "Fintech growth ($2.1B market) expanding access to underserved segments",
            "AI and alternative data improving credit scoring accuracy",
            "Mobile money integration reducing transaction costs",
            "Embedded finance bringing credit to point of need",
            "Digital identity systems enhancing KYC and reducing fraud"
        ],
        21
    )

    # Slide 22: Enhanced Infrastructure
    add_content_slide(
        "Opportunity 2: Strengthening Credit Infrastructure",
        [
            "Expand credit bureau coverage and reporting comprehensiveness",
            "Promote NCR awareness and movable collateral utilization",
            "Integrate digital lending data into formal credit reporting systems",
            "Develop centralized credit scoring models for underserved segments",
            "Establish specialized commercial courts for faster loan recovery",
            "Harmonize credit data standards across all lender categories"
        ],
        22
    )

    # Slide 23: Targeted Interventions
    add_content_slide(
        "Opportunity 3: Targeted Sector Interventions",
        [
            "Scale up ACGSF and SMECGS to reach more beneficiaries",
            "Develop sector-specific credit products (agriculture, manufacturing, tech)",
            "Increase DBN capitalization for expanded MSME on-lending",
            "Create gender-focused credit programs leveraging DBN's 66% women reach",
            "Support value chain financing to de-risk lending",
            "Subsidize interest rates for strategic sectors"
        ],
        23
    )

    # Slide 24: Regulatory Improvements
    add_content_slide(
        "Opportunity 4: Regulatory and Legal Reforms",
        [
            "Strengthen creditor rights and enforcement mechanisms",
            "Fast-track insolvency proceedings for loan recovery",
            "Standardize digital lending regulations across all agencies",
            "Mandate comprehensive credit bureau reporting with penalties",
            "Simplify movable asset registration and enforcement processes",
            "Harmonize regulatory frameworks across CBN, FCCPC, SEC"
        ],
        24
    )

    # Slide 25: Capacity Building
    add_content_slide(
        "Opportunity 5: Capacity Building and Awareness",
        [
            "Financial literacy programs for borrowers on credit management",
            "Train lenders on alternative credit assessment methodologies",
            "Promote awareness of guarantee schemes and NCR among SMEs",
            "Build capacity of MFBs and community banks for risk management",
            "Develop standardized business documentation templates for SMEs",
            "Leverage technology for scalable financial education"
        ],
        25
    )

    # Slide 26: Key Recommendations
    add_content_slide(
        "Key Recommendations",
        [
            "Accelerate full implementation of Open Banking for better credit decisions",
            "Mandate universal credit bureau reporting with regulatory teeth",
            "Increase capitalization of guarantee schemes and DFIs",
            "Fast-track legal reforms for creditor rights and loan recovery",
            "Promote alternative collateral through policy incentives for lenders",
            "Establish national financial literacy program with credit focus",
            "Create conducive macroeconomic environment for sustainable lending"
        ],
        26
    )

    # Slide 27: Conclusion
    add_content_slide(
        "Conclusion",
        [
            "Credit facilitation is critical for Nigeria's economic transformation",
            "Significant progress made in infrastructure and regulation (2025-2026)",
            "Persistent challenges require coordinated multi-stakeholder response",
            "Digital innovation presents transformative opportunities",
            "Success requires alignment of policy, regulation, and market practice",
            "The $32.2B MSME credit gap represents both challenge and opportunity"
        ],
        27
    )

    # Slide 28: Thank You / Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.fill.background()

    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(54)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = RGBColor(255, 255, 255)
    thanks_para.alignment = PP_ALIGN.CENTER

    # Q&A text
    qa_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    qa_frame = qa_box.text_frame
    qa_frame.text = "Questions & Discussion"
    qa_para = qa_frame.paragraphs[0]
    qa_para.font.size = Pt(32)
    qa_para.font.color.rgb = RGBColor(200, 200, 200)
    qa_para.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save('/home/user/Kenna/Credit_Facilitation_Nigeria.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
